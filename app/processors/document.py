"""
Document processing for Word, Excel, and PowerPoint files
"""
import os
from typing import Dict, Any, List


def extract_text_from_docx(docx_path: str) -> str:
    """
    Extract text from Word document (.docx).
    
    Args:
        docx_path: Path to Word file
    
    Returns:
        Extracted text content
    """
    try:
        from docx import Document
        
        doc = Document(docx_path)
        text_parts = []
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        return "\n\n".join(text_parts)
        
    except Exception as e:
        return f"Error extracting text from Word document: {str(e)}"


def extract_text_from_xlsx(xlsx_path: str) -> str:
    """
    Extract text/data from Excel file (.xlsx).
    
    Args:
        xlsx_path: Path to Excel file
    
    Returns:
        Extracted text content (all sheets)
    """
    try:
        from openpyxl import load_workbook
        
        wb = load_workbook(xlsx_path, data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===")
            
            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        row_values.append(str(cell.value))
                if row_values:
                    text_parts.append(" | ".join(row_values))
        
        wb.close()
        return "\n".join(text_parts)
        
    except Exception as e:
        return f"Error extracting text from Excel file: {str(e)}"


def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extract text from PowerPoint file (.pptx).
    
    Args:
        pptx_path: Path to PowerPoint file
    
    Returns:
        Extracted text content (all slides)
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"=== Slide {slide_num} ==="]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract text from tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            if len(slide_text) > 1:  # More than just the slide header
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
        
    except Exception as e:
        return f"Error extracting text from PowerPoint file: {str(e)}"


def extract_document_metadata(file_path: str, file_type: str) -> Dict[str, Any]:
    """
    Extract metadata from document files.
    
    Args:
        file_path: Path to document file
        file_type: Type of document (docx, xlsx, pptx)
    
    Returns:
        Dictionary with document metadata
    """
    metadata = {
        "file_path": file_path,
        "file_name": os.path.basename(file_path),
        "file_size": os.path.getsize(file_path) if os.path.exists(file_path) else 0,
        "type": file_type
    }
    
    try:
        if file_type == "docx":
            from docx import Document
            doc = Document(file_path)
            core_props = doc.core_properties
            metadata.update({
                "author": core_props.author or "",
                "title": core_props.title or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables)
            })
            
        elif file_type == "xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(file_path, data_only=True)
            metadata.update({
                "sheet_count": len(wb.sheetnames),
                "sheet_names": wb.sheetnames
            })
            wb.close()
            
        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            metadata.update({
                "slide_count": len(prs.slides)
            })
            
    except Exception as e:
        metadata["metadata_error"] = str(e)
    
    return metadata


def process_document(file_path: str, file_type: str) -> Dict[str, Any]:
    """
    Full document processing pipeline.
    
    Args:
        file_path: Path to document file
        file_type: Type of document (docx, xlsx, pptx)
    
    Returns:
        Dictionary with extracted text and metadata
    """
    result = {
        "file_path": file_path,
        "file_name": os.path.basename(file_path),
        "type": file_type
    }
    
    # Extract metadata
    metadata = extract_document_metadata(file_path, file_type)
    result["metadata"] = metadata
    
    # Extract text based on type
    if file_type == "docx":
        result["text"] = extract_text_from_docx(file_path)
    elif file_type == "xlsx":
        result["text"] = extract_text_from_xlsx(file_path)
    elif file_type == "pptx":
        result["text"] = extract_text_from_pptx(file_path)
    else:
        result["text"] = ""
        result["error"] = f"Unsupported document type: {file_type}"
    
    return result
